import os
import requests
from dotenv import load_dotenv
from PIL import Image
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import pandas as pd
import tempfile

# Load environment variables
load_dotenv()
MISTRAL_API_KEY = os.getenv('MISTRAL_API_KEY')
MISTRAL_API_URL = 'https://api.mistral.ai/v1/ocr'

def process_with_mistral_ocr(file_path, file_type):
    headers = {
        'Authorization': f'Bearer {MISTRAL_API_KEY}',
        'Content-Type': 'application/octet-stream'
    }

    try:
        with open(file_path, 'rb') as file:
            files = {'file': (os.path.basename(file_path), file)}
            params = {'type': file_type}

            response = requests.post(
                MISTRAL_API_URL,
                headers=headers,
                files=files,
                params=params
            )

            if response.status_code == 200:
                return response.json().get('text', '')
            else:
                print(f"Error with Mistral OCR: {response.status_code} - {response.text}")
                return None
    except Exception as e:
        print(f"Error processing file with Mistral OCR: {e}")
        return None

def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = '\n'.join([page.extract_text() or '' for page in reader.pages])
        if text.strip():
            return text
    except:
        pass

    try:
        images = convert_from_path(pdf_path)
        extracted_text = []

        for image in images:
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                image_path = temp_file.name
                image.save(image_path, 'JPEG')

            text = process_with_mistral_ocr(image_path, 'image')
            if text:
                extracted_text.append(text)

            os.unlink(image_path)

        return '\n'.join(extracted_text)
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_image(image_path):
    return process_with_mistral_ocr(image_path, 'image')

def extract_text_from_word(docx_path):
    try:
        doc = Document(docx_path)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_part = rel.target_part
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                    image_path = temp_file.name
                    with open(image_path, 'wb') as f:
                        f.write(image_part.blob)

                    image_text = extract_text_from_image(image_path)
                    if image_text:
                        text += f"\n[IMAGE CONTENT]\n{image_text}\n"

                os.unlink(image_path)

        return text
    except Exception as e:
        print(f"Error extracting text from Word document: {e}")
        return None

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

                if shape.shape_type == 13:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                        image_path = temp_file.name
                        with open(image_path, 'wb') as f:
                            f.write(shape.image.blob)

                        image_text = extract_text_from_image(image_path)
                        if image_text:
                            text.append(f"[IMAGE CONTENT]\n{image_text}")

                    os.unlink(image_path)

        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return None

def extract_text_from_excel(excel_path):
    try:
        df = pd.read_excel(excel_path, sheet_name=None)
        text = []
        for sheet_name, sheet_data in df.items():
            text.append(f"=== Sheet: {sheet_name} ===")
            text.append(sheet_data.to_string())
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return None

def process_file(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return None

    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ('.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff'):
        return extract_text_from_image(file_path)
    elif file_ext in ('.docx', '.doc'):
        return extract_text_from_word(file_path)
    elif file_ext in ('.pptx', '.ppt'):
        return extract_text_from_pptx(file_path)
    elif file_ext in ('.xlsx', '.xls', '.csv'):
        return extract_text_from_excel(file_path)
    else:
        print(f"Unsupported file type: {file_ext}")
        return None

if __name__ == "__main__":
    file_path = input("Masukkan path file yang ingin diproses: ")
    extracted_text = process_file(file_path)

    if extracted_text:
        print("\n=== Teks yang Diekstrak ===")
        print(extracted_text[:2000] + "..." if len(extracted_text) > 2000 else extracted_text)

        output_path = os.path.splitext(file_path)[0] + '_extracted.txt'
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
        print(f"\nHasil disimpan ke: {output_path}")
    else:
        print("Gagal mengekstrak teks dari file.")
